"""
Samvaadhika - AI Processing Pipeline
Handles: language detection → ASR → MT → TTS / subtitles / document re-assembly.

All models run locally (CPU). On first use each model is loaded once and cached
in memory for the lifetime of the process.

Stubs are provided so the app runs end-to-end even before the heavy AI models
are downloaded — each stub logs a clear message and returns a placeholder result.
"""
import hashlib
import logging
import os
import shutil
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional, Tuple

from app.config import (
    CACHE_DIR, OUTPUTS_DIR, UPLOADS_DIR, MODELS_DIR,
    WHISPER_MODEL_SIZE, WHISPER_DEVICE, WHISPER_COMPUTE_TYPE,
    TESSERACT_LANGUAGES, SUPPORTED_LANGUAGES, BASE_DIR,
)

logger = logging.getLogger("samvaadhika.pipeline")

# ---------------------------------------------------------------------------
# Ensure ffmpeg is on PATH (winget installs to a deep location)
# ---------------------------------------------------------------------------
_FFMPEG_SEARCH_DIRS = [
    Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "WinGet" / "Packages",
    Path(os.environ.get("ProgramFiles", "C:/Program Files")) / "ffmpeg" / "bin",
    Path(os.environ.get("ProgramFiles(x86)", "C:/Program Files (x86)")) / "ffmpeg" / "bin",
]

def _ensure_ffmpeg_on_path():
    """Find ffmpeg installed by winget (or other locations) and add to PATH."""
    if shutil.which("ffmpeg"):
        return  # already available
    for search_root in _FFMPEG_SEARCH_DIRS:
        if not search_root.exists():
            continue
        for ffmpeg_exe in search_root.rglob("ffmpeg.exe"):
            bin_dir = str(ffmpeg_exe.parent)
            os.environ["PATH"] = bin_dir + os.pathsep + os.environ.get("PATH", "")
            logger.info(f"Added ffmpeg to PATH from: {bin_dir}")
            return
    logger.warning("ffmpeg not found in any known location. Video/audio processing may fail.")

_ensure_ffmpeg_on_path()


# ---------------------------------------------------------------------------
# HuggingFace cache directory resolver
# ---------------------------------------------------------------------------

def _resolve_hf_cache(base_dir: Path) -> Path:
    """
    Resolve a HuggingFace cache directory to the actual model snapshot path.

    HuggingFace downloads create a structure like:
        models/model-name/models--org--model-name/snapshots/<hash>/
    This function finds the actual snapshot directory containing model files.
    If base_dir itself contains model files directly, returns base_dir as-is.
    """
    # If the directory directly contains model files, use it as-is
    direct_indicators = ["config.json", "model.bin", "model.safetensors",
                         "tokenizer.json", "tokenizer_config.json"]
    for indicator in direct_indicators:
        if (base_dir / indicator).exists():
            return base_dir

    # Look for HuggingFace cache structure: models--*/snapshots/*/
    for models_dir in base_dir.glob("models--*"):
        snapshots_dir = models_dir / "snapshots"
        if snapshots_dir.exists():
            # Get the latest snapshot (usually only one)
            snapshots = sorted(snapshots_dir.iterdir(), key=lambda p: p.stat().st_mtime, reverse=True)
            for snap in snapshots:
                if snap.is_dir() and any(snap.iterdir()):
                    logger.info(f"Resolved HF cache: {base_dir.name} → {snap}")
                    return snap

    # Fallback: return the original directory
    return base_dir


def _find_model_dir(name: str, *alt_names: str) -> Optional[Path]:
    """
    Find a model directory under MODELS_DIR, trying multiple name variants.
    Returns the resolved path (handling HF cache structure) or None.
    """
    candidates = [name] + list(alt_names)
    for candidate in candidates:
        model_dir = MODELS_DIR / candidate
        if model_dir.exists():
            return _resolve_hf_cache(model_dir)
    return None


# ---------------------------------------------------------------------------
# Lazy model singletons
# ---------------------------------------------------------------------------
_whisper_model = None
_lang_detector = None


def _get_whisper():
    global _whisper_model
    if _whisper_model is None:
        try:
            from faster_whisper import WhisperModel

            # Try local model directory first, then fall back to model size string
            local_model = _find_model_dir(
                f"faster-whisper-{WHISPER_MODEL_SIZE}",
                f"whisper-{WHISPER_MODEL_SIZE}",
                "faster-whisper",
            )
            model_source = str(local_model) if local_model else WHISPER_MODEL_SIZE

            _whisper_model = WhisperModel(
                model_source,
                device=WHISPER_DEVICE,
                compute_type=WHISPER_COMPUTE_TYPE,
            )
            logger.info(f"Whisper model loaded from: {model_source}")
        except Exception as e:
            logger.warning(f"faster-whisper not available: {e}. Using stub ASR.")
            _whisper_model = "stub"
    return _whisper_model


def _get_lang_detector():
    global _lang_detector
    if _lang_detector is None:
        try:
            import fasttext
            # Try both .ftz (compressed) and .bin (full) variants
            model_path = None
            for fname in ["lid.176.ftz", "lid.176.bin"]:
                candidate = MODELS_DIR / fname
                if candidate.exists():
                    model_path = candidate
                    break

            if model_path:
                _lang_detector = fasttext.load_model(str(model_path))
                logger.info(f"fastText language detector loaded from {model_path.name}.")
            else:
                logger.warning("fastText model not found (tried lid.176.ftz, lid.176.bin) — using stub.")
                _lang_detector = "stub"
        except Exception as e:
            logger.warning(f"fasttext not available: {e}. Using stub language detector.")
            _lang_detector = "stub"
    return _lang_detector


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()


def _ffmpeg_available() -> bool:
    return shutil.which("ffmpeg") is not None


def _tesseract_available() -> bool:
    return shutil.which("tesseract") is not None


# ---------------------------------------------------------------------------
# Language detection
# ---------------------------------------------------------------------------

def detect_language(text: str) -> str:
    """Return ISO 639-1 code: 'en', 'hi', or 'mr'. Falls back to 'en'."""
    detector = _get_lang_detector()
    if detector == "stub" or not text.strip():
        return "en"
    try:
        # Avoid fasttext's NumPy-2-incompatible Python predict wrapper.
        predictions = detector.f.predict(text.replace("\n", " "), 1, 0.0, "strict")
        label = predictions[0][1].replace("__label__", "")
        # fastText uses 'hi' and 'mr' directly
        if label in SUPPORTED_LANGUAGES:
            return label
        return "en"
    except Exception as e:
        logger.warning(f"Language detection failed: {e}")
        return "en"


# ---------------------------------------------------------------------------
# Translation (MT)
# ---------------------------------------------------------------------------

def translate_text(text: str, source_lang: str, target_lang: str) -> Tuple[str, float]:
    """
    Translate text using IndicTrans2 (preferred) or argostranslate (fallback).
    Returns (translated_text, confidence_score 0-1).
    """
    if source_lang == target_lang:
        logger.info(f"translate_text: source==target ({source_lang}); skipping MT")
        return text, 1.0

    # Try IndicTrans2 first
    try:
        res = _translate_indictrans2(text, source_lang, target_lang)
        logger.info("translate_text: used IndicTrans2")
        return res
    except Exception as e:
        logger.warning(f"IndicTrans2 unavailable ({e}), trying argostranslate fallback.")

    # Argostranslate fallback
    try:
        res = _translate_argos(text, source_lang, target_lang)
        logger.info("translate_text: used ArgosTranslate fallback")
        return res
    except Exception as e:
        logger.warning(f"argostranslate unavailable ({e}). Returning stub translation.")

    # Final stub — clearly marked so reviewers know it's a placeholder
    return f"[TRANSLATION STUB: {source_lang}→{target_lang}] {text}", 0.0


def _translate_indictrans2(text: str, src: str, tgt: str) -> Tuple[str, float]:
    """
    IndicTrans2 distilled model with direction-based routing.

    Supported models:
    - indictrans2-en-indic-dist-200M: English → Hindi, Marathi (available)
    - indictrans2-indic-en-dist-200M: Hindi, Marathi → English (gated, need access request)
    - indictrans2-indic-indic-dist-320M: Hindi ↔ Marathi (gated, need access request)

    Returns (translated_text, confidence_score).
    Confidence reflects model quality and direction support.
    """
    from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
    import torch

    # Direction-to-model mapping
    direction_map = {
        ('en', 'hi'): ('indictrans2-en-indic-dist-200M', 0.88),
        ('en', 'mr'): ('indictrans2-en-indic-dist-200M', 0.88),
        ('hi', 'en'): ('indictrans2-indic-en-dist-200M', 0.85),
        ('mr', 'en'): ('indictrans2-indic-en-dist-200M', 0.85),
        ('hi', 'mr'): ('indictrans2-indic-indic-dist-320M', 0.82),
        ('mr', 'hi'): ('indictrans2-indic-indic-dist-320M', 0.82),
    }

    # Check if direction is supported
    direction = (src, tgt)
    if direction not in direction_map:
        raise ValueError(
            f"IndicTrans2 does not support {src}→{tgt} translation. "
            f"Supported directions: en↔hi, en↔mr, hi↔mr"
        )

    model_name, confidence = direction_map[direction]

    # Find the model directory
    model_dir = _find_model_dir(model_name, "indictrans2")
    if model_dir is None:
        raise FileNotFoundError(
            f"IndicTrans2 model '{model_name}' not found at models/{model_name}/. "
            f"For {src}→{tgt}, the gated repository access may be required. "
            f"Visit: https://huggingface.co/ai4bharat/{model_name} to request access."
        )

    # Language code mapping for IndicTrans2
    lang_map = {"en": "eng_Latn", "hi": "hin_Deva", "mr": "mar_Deva"}
    src_code = lang_map.get(src, "eng_Latn")
    tgt_code = lang_map.get(tgt, "hin_Deva")

    tokenizer = AutoTokenizer.from_pretrained(str(model_dir), trust_remote_code=True)
    model = AutoModelForSeq2SeqLM.from_pretrained(str(model_dir), trust_remote_code=True)
    model.eval()

    # IndicTrans2 custom tokenizer expects: "src_lang tgt_lang actual_text"
    tagged_text = f"{src_code} {tgt_code} {text}"
    inputs = tokenizer(tagged_text, return_tensors="pt", padding=True, truncation=True, max_length=512)

    with torch.no_grad():
        outputs = model.generate(**inputs, max_length=512, num_beams=1, use_cache=False)

    # Switch tokenizer to target mode for decoding, then restore
    tokenizer._switch_to_target_mode()
    translated = tokenizer.decode(outputs[0], skip_special_tokens=True)
    tokenizer._switch_to_input_mode()

    return translated, confidence


def _translate_argos(text: str, src: str, tgt: str) -> Tuple[str, float]:
    """Argostranslate CPU fallback — limited Indic support but works offline."""
    import argostranslate.package
    import argostranslate.translate

    installed = argostranslate.translate.get_installed_languages()
    src_lang = next((l for l in installed if l.code == src), None)
    tgt_lang = next((l for l in installed if l.code == tgt), None)

    if src_lang is None or tgt_lang is None:
        raise RuntimeError(f"Argostranslate language pair {src}→{tgt} not installed.")

    translation = src_lang.get_translation(tgt_lang)
    result = translation.translate(text)
    return result, 0.6


# ---------------------------------------------------------------------------
# Glossary application
# ---------------------------------------------------------------------------

def apply_glossary(text: str, source_lang: str, target_lang: str, db) -> str:
    """Replace known domain terms in the translated text using the glossary table."""
    try:
        from app.models import GlossaryEntry
        entries = (
            db.query(GlossaryEntry)
            .filter(
                GlossaryEntry.source_language == source_lang,
                GlossaryEntry.target_language == target_lang,
            )
            .all()
        )
        for entry in entries:
            text = text.replace(entry.source_term, entry.target_term)
    except Exception as e:
        logger.warning(f"Glossary application failed: {e}")
    return text


# ---------------------------------------------------------------------------
# ASR — Speech to Text
# ---------------------------------------------------------------------------

def transcribe_audio(audio_path: Path, language: Optional[str] = None) -> Tuple[list, str]:
    """
    Transcribe audio file. Returns (segments, detected_language).
    Each segment: {"start": float, "end": float, "text": str}
    """
    model = _get_whisper()
    if model == "stub":
        logger.warning("ASR stub: returning placeholder transcript.")
        return [{"start": 0.0, "end": 5.0, "text": "[ASR not available — install faster-whisper]"}], "en"

    try:
        # First attempt with VAD enabled (faster, skips silence)
        segments_iter, info = model.transcribe(
            str(audio_path),
            language=language,
            beam_size=5,
            vad_filter=True,
        )
        segments = [
            {"start": s.start, "end": s.end, "text": s.text.strip()}
            for s in segments_iter
        ]

        # If VAD produced very little speech for a long input, retry without VAD
        try:
            import wave
            duration = None
            if audio_path.suffix.lower() == ".wav":
                try:
                    with wave.open(str(audio_path), 'rb') as wf:
                        duration = wf.getnframes() / float(wf.getframerate())
                except Exception:
                    duration = None

            total_speech = sum((s["end"] - s["start"]) for s in segments) if segments else 0.0
            if duration and total_speech < max(1.0, duration * 0.15):
                logger.info(f"ASR VAD produced only {total_speech:.1f}s speech from {duration:.1f}s audio; retrying without VAD.")
                segments_iter, info = model.transcribe(
                    str(audio_path),
                    language=language,
                    beam_size=5,
                    vad_filter=False,
                )
                segments = [
                    {"start": s.start, "end": s.end, "text": s.text.strip()}
                    for s in segments_iter
                ]
        except Exception:
            # best-effort duration check; continue with whatever segments we have
            pass
        # Debug: log brief summary of segments
        try:
            if segments:
                logger.info(f"ASR: {len(segments)} segments. First segment: '{segments[0]['text'][:200]}'")
            else:
                logger.info("ASR: no segments produced.")
        except Exception:
            pass
        return segments, info.language
    except Exception as e:
        logger.error(f"ASR transcription failed: {e}")
        raise


# ---------------------------------------------------------------------------
# TTS — Text to Speech
# ---------------------------------------------------------------------------

def synthesize_speech(text: str, language: str, output_path: Path) -> bool:
    """
    Generate speech audio from text using Indic Parler-TTS (preferred)
    or pyttsx3 stub fallback.
    Returns True on success.
    """
    try:
        return _tts_parler(text, language, output_path)
    except Exception as e:
        logger.warning(f"Parler-TTS unavailable ({e}), trying pyttsx3 stub.")
        logger.debug("Parler-TTS exception details:", exc_info=True)

    try:
        return _tts_pyttsx3(text, language, output_path)
    except Exception as e:
        logger.warning(f"pyttsx3 unavailable ({e}). TTS skipped.")
        return False


def _tts_parler(text: str, language: str, output_path: Path) -> bool:
    """AI4Bharat Indic Parler-TTS — Apache-2.0 licensed."""
    import torch
    # Workaround: disable TorchScript compilation (monkeypatch torch.jit.script/trace)
    # Apply monkeypatch BEFORE importing `parler_tts` so any TorchScript attempts
    # during module import are no-ops.
    _orig_jit_script = getattr(torch.jit, "script", None)
    _orig_jit_trace = getattr(torch.jit, "trace", None)
    def _noop_jit(x, *a, **k):
        return x
    try:
        if _orig_jit_script is not None:
            torch.jit.script = _noop_jit
        if _orig_jit_trace is not None:
            torch.jit.trace = _noop_jit
    except Exception:
        pass

    # Now import parler_tts and related modules (monkeypatched)
    try:
        from parler_tts import ParlerTTSForConditionalGeneration
        from transformers import AutoTokenizer
        import soundfile as sf
    except Exception as e:
        # Restore torch.jit before re-raising
        try:
            if _orig_jit_script is not None:
                torch.jit.script = _orig_jit_script
            if _orig_jit_trace is not None:
                torch.jit.trace = _orig_jit_trace
        except Exception:
            pass
        raise

    model_dir = _find_model_dir("indic-parler-tts", "parler-tts")
    if model_dir is None:
        raise FileNotFoundError(
            "Parler-TTS model not found. Expected at models/indic-parler-tts/"
        )

    tokenizer = AutoTokenizer.from_pretrained(str(model_dir))
    model = ParlerTTSForConditionalGeneration.from_pretrained(str(model_dir))
    model.eval()

    description = "A female speaker delivers a clear, natural voice."
    input_ids = tokenizer(description, return_tensors="pt").input_ids
    prompt_ids = tokenizer(text, return_tensors="pt").input_ids

    with torch.no_grad():
        generation = model.generate(input_ids=input_ids, prompt_input_ids=prompt_ids)

    try:
        audio = generation.cpu().numpy().squeeze()
        sf.write(str(output_path), audio, model.config.sampling_rate)
        return True
    finally:
        # Restore original torch.jit functions
        try:
            if _orig_jit_script is not None:
                torch.jit.script = _orig_jit_script
            if _orig_jit_trace is not None:
                torch.jit.trace = _orig_jit_trace
        except Exception:
            pass


def _tts_pyttsx3(text: str, language: str, output_path: Path) -> bool:
    """pyttsx3 system TTS stub — English only, for dev/demo."""
    import pyttsx3
    engine = pyttsx3.init()
    engine.save_to_file(text, str(output_path))
    engine.runAndWait()
    # Verify output file was written and has non-trivial size
    try:
        if output_path.exists() and output_path.stat().st_size > 1024:
            logger.info(f"pyttsx3 produced audio file: {output_path} ({output_path.stat().st_size} bytes)")
            return True
        else:
            logger.warning(f"pyttsx3 produced empty or tiny audio file: {output_path} ({output_path.stat().st_size if output_path.exists() else 0} bytes)")
            return False
    except Exception as e:
        logger.warning(f"pyttsx3 verification failed: {e}")
        return False


# ---------------------------------------------------------------------------
# Subtitle generation
# ---------------------------------------------------------------------------

def generate_subtitles(segments: list, translated_segments: list, output_path: Path) -> bool:
    """Generate SRT subtitle file from timed segments."""
    try:
        import pysubs2
        subs = pysubs2.SSAFile()
        for orig, trans in zip(segments, translated_segments):
            event = pysubs2.SSAEvent(
                start=pysubs2.make_time(s=orig["start"]),
                end=pysubs2.make_time(s=orig["end"]),
                text=trans["text"],
            )
            subs.append(event)
        subs.save(str(output_path))
        return True
    except Exception as e:
        logger.error(f"Subtitle generation failed: {e}")
        return False


# ---------------------------------------------------------------------------
# Document translation
# ---------------------------------------------------------------------------

def translate_docx(input_path: Path, output_path: Path, source_lang: str, target_lang: str, db) -> bool:
    """Translate a DOCX file, preserving formatting."""
    try:
        from docx import Document
        doc = Document(str(input_path))
        for para in doc.paragraphs:
            if para.text.strip():
                translated, _ = translate_text(para.text, source_lang, target_lang)
                translated = apply_glossary(translated, source_lang, target_lang, db)
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = translated
                else:
                    para.text = translated
        # Tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if para.text.strip():
                            translated, _ = translate_text(para.text, source_lang, target_lang)
                            translated = apply_glossary(translated, source_lang, target_lang, db)
                            for run in para.runs:
                                run.text = ""
                            if para.runs:
                                para.runs[0].text = translated
        doc.save(str(output_path))
        return True
    except Exception as e:
        logger.error(f"DOCX translation failed: {e}")
        raise


def translate_pptx(input_path: Path, output_path: Path, source_lang: str, target_lang: str, db) -> bool:
    """Translate a PPTX file, preserving slide layout."""
    try:
        from pptx import Presentation
        prs = Presentation(str(input_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                translated, _ = translate_text(run.text, source_lang, target_lang)
                                translated = apply_glossary(translated, source_lang, target_lang, db)
                                run.text = translated
        prs.save(str(output_path))
        return True
    except Exception as e:
        logger.error(f"PPTX translation failed: {e}")
        raise


def _insert_pdf_text(page, rect, text: str, font_path: Optional[Path], fontsize: float) -> bool:
    """Render translated text as an image so Devanagari shaping works reliably."""
    from io import BytesIO
    import fitz
    from PIL import Image, ImageDraw, ImageFont

    scale = 3
    width = max(1, int(rect.width * scale))
    height = max(1, int(rect.height * scale))
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font = None
    if font_path:
        try:
            font = ImageFont.truetype(str(font_path), max(8, int(fontsize * scale)), index=0)
        except Exception:
            font = None
    if font is None:
        font = ImageFont.load_default()

    words = text.split()
    lines = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if current and draw.textbbox((0, 0), candidate, font=font)[2] > width - 8:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    if not lines:
        return False

    line_height = max(1, int(fontsize * scale * 1.15))
    # Compute required image height for all lines and expand image if needed
    padding_px = 6
    required_height = len(lines) * line_height + padding_px
    if required_height > height:
        # create a taller image and copy existing white background
        new_image = Image.new("RGB", (width, required_height), "white")
        new_image.paste(image, (0, 0))
        image = new_image
        draw = ImageDraw.Draw(image)

    draw.multiline_text((4, 1), "\n".join(lines), font=font, fill="black", spacing=0)
    stream = BytesIO()
    image.save(stream, format="PNG")
    # If we expanded the image height, map the image back to a taller rect on the page
    if image.height != height:
        # compute new rect height in page coordinate space
        new_height_pts = image.height / scale
        new_rect = fitz.Rect(rect.x0, rect.y0, rect.x1, rect.y0 + new_height_pts)
        # ensure we don't overflow the page bottom
        page_bottom = page.rect.y1
        if new_rect.y1 > page_bottom:
            shift_up = new_rect.y1 - page_bottom
            new_rect = fitz.Rect(new_rect.x0, max(rect.y0 - shift_up, page.rect.y0), new_rect.x1, page_bottom)
        page.insert_image(new_rect, stream=stream.getvalue(), overlay=True)
    else:
        page.insert_image(rect, stream=stream.getvalue(), overlay=True)
    return True


def translate_pdf(input_path: Path, output_path: Path, source_lang: str, target_lang: str, db) -> Tuple[bool, str]:
    """
    Translate a PDF while preserving its page geometry and table/grid layout.
    Returns (success, notes).
    """
    try:
        import fitz
        import pdfplumber
        output_path.parent.mkdir(parents=True, exist_ok=True)
        formatted_pages = 0
        fallback_pages = 0

        # Prefer a system or deployment-provided Devanagari font for Hindi/Marathi.
        font_candidates = [
            os.environ.get("SAMVAADHIKA_DEVANAGARI_FONT", ""),
            "/System/Library/Fonts/Kohinoor.ttc",
            "/System/Library/Fonts/Supplemental/Devanagari Sangam MN.ttc",
            "C:/Windows/Fonts/mangal.ttf",
        ]
        # Also prefer any fonts bundled with the application under BASE_DIR/fonts
        bundled_fonts_dir = BASE_DIR / "fonts"
        if bundled_fonts_dir.exists():
            # Prefer explicit known names first
            for candidate_name in ("NotoSansDevanagari-Regular.ttf", "mangal.ttf", "NotoSansDevanagari.ttc"):
                candidate = bundled_fonts_dir / candidate_name
                if candidate.exists():
                    font_candidates.insert(0, str(candidate))
            # Otherwise, add any ttf/ttc in the bundled fonts folder
            for f in bundled_fonts_dir.iterdir():
                if f.suffix.lower() in (".ttf", ".ttc"):
                    font_candidates.append(str(f))

        font_path = next((Path(p) for p in font_candidates if p and Path(p).exists()), None)
        if font_path:
            logger.info(f"Using Devanagari font: {font_path}")
        formatted_doc = fitz.open(str(input_path))

        with pdfplumber.open(str(input_path)) as pdf:
            for page_number, source_page in enumerate(pdf.pages):
                output_page = formatted_doc[page_number]
                table_cells = []
                for table in source_page.find_tables():
                    table_cells.extend(cell for row in table.rows for cell in row.cells if cell)

                if table_cells:
                    translated_any = False
                    for x0, top, x1, bottom in table_cells:
                        cell_rect = fitz.Rect(x0, top, x1, bottom)
                        cell_text = source_page.crop((x0, top, x1, bottom)).extract_text() or ""
                        if not cell_text.strip():
                            continue
                        translated, _ = translate_text(cell_text, source_lang, target_lang)
                        translated = apply_glossary(translated, source_lang, target_lang, db)
                        words = output_page.get_text("words", clip=cell_rect)
                        if words:
                            text_rect = fitz.Rect(
                                min(word[0] for word in words) - 1,
                                min(word[1] for word in words) - 1,
                                max(word[2] for word in words) + 1,
                                max(word[3] for word in words) + 1,
                            ) & cell_rect
                            output_page.draw_rect(text_rect, color=None, fill=(1, 1, 1), overlay=True)
                        translated_any = _insert_pdf_text(
                            output_page,
                            fitz.Rect(cell_rect.x0 + 2, cell_rect.y0 + 1, cell_rect.x1 - 2, cell_rect.y1 - 1),
                            translated,
                            font_path,
                            max(5, min(10, (bottom - top) * 0.42)),
                        ) or translated_any
                    if translated_any:
                        formatted_pages += 1
                        continue

                # For non-table text-native pages, preserve each text line's position.
                words = source_page.extract_words(keep_blank_chars=True, use_text_flow=True)
                lines = {}
                for word in words:
                    key = (round(word["top"], 1), round(word["bottom"], 1))
                    lines.setdefault(key, []).append(word)
                translated_any = False
                for (top, bottom), line_words in lines.items():
                    line_words.sort(key=lambda word: word["x0"])
                    original = " ".join(word["text"] for word in line_words).strip()
                    if not original:
                        continue
                    translated, _ = translate_text(original, source_lang, target_lang)
                    translated = apply_glossary(translated, source_lang, target_lang, db)
                    text_rect = fitz.Rect(
                        min(word["x0"] for word in line_words), top,
                        max(word["x1"] for word in line_words), bottom,
                    )
                    output_page.draw_rect(text_rect, color=None, fill=(1, 1, 1), overlay=True)
                    translated_any = _insert_pdf_text(
                        output_page,
                        text_rect,
                        translated,
                        font_path,
                        max(5, min(11, bottom - top)),
                    ) or translated_any
                if translated_any:
                    formatted_pages += 1
                else:
                    fallback_pages += 1

        if fallback_pages:
            logger.warning("%d PDF page(s) had no coordinate-aware text; layout may be incomplete.", fallback_pages)
        formatted_doc.save(str(output_path), garbage=4, deflate=True)
        formatted_doc.close()
        return True, f"PDF translated with layout preservation on {formatted_pages} page(s)."
    except Exception as e:
        logger.error(f"PDF translation failed: {e}")
        raise


# ---------------------------------------------------------------------------
# Audio/Video extraction helpers
# ---------------------------------------------------------------------------

def extract_audio_from_video(video_path: Path, audio_path: Path) -> bool:
    """Use ffmpeg to extract audio track from video."""
    if not _ffmpeg_available():
        raise RuntimeError("ffmpeg not found. Install ffmpeg and add to PATH.")
    cmd = [
        "ffmpeg", "-y", "-i", str(video_path),
        "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1",
        str(audio_path),
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=300)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg audio extraction failed: {result.stderr.decode()}")
    return True


def normalize_audio(input_path: Path, output_path: Path) -> bool:
    """Normalize audio to 16kHz mono WAV for Whisper."""
    if not _ffmpeg_available():
        shutil.copy(input_path, output_path)
        return True
    cmd = [
        "ffmpeg", "-y", "-i", str(input_path),
        "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1",
        str(output_path),
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=300)
    return result.returncode == 0
